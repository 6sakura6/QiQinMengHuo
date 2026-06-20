#!/usr/bin/env python3
"""
《七擒孟获》游戏介绍PPT — 优化版生成脚本
基于 v1.1 游戏文档内容深度优化
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ============================================================
# 配色方案：像素三国主题
# ============================================================
# 主色
C_PRIMARY     = RGBColor(0xDC, 0x26, 0x26)  # 朱砂红
C_DARK        = RGBColor(0x0F, 0x17, 0x2A)   # 深色底
C_ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x1E)   # 古铜金
C_ACCENT_GREEN= RGBColor(0x22, 0xC5, 0x5E)   # 强调绿
C_INK_GREEN   = RGBColor(0x1A, 0x3C, 0x2A)   # 墨绿
C_PURPLE      = RGBColor(0x7C, 0x3A, 0xED)   # 暗紫
C_DEEP_BLUE   = RGBColor(0x1E, 0x3A, 0x5F)   # 深蓝

# 中性色
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY  = RGBColor(0xF1, 0xF0, 0xEB)   # 米白纸色
C_MID_GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
C_DARK_GRAY   = RGBColor(0x37, 0x41, 0x51)
C_TEXT_DARK   = RGBColor(0x1F, 0x29, 0x37)
C_TEXT_BODY   = RGBColor(0x4B, 0x55, 0x63)

# 关卡色调
LEVEL_COLORS = [
    RGBColor(0x4A, 0x90, 0x3C),  # 1 翠绿
    RGBColor(0x3B, 0x82, 0xC4),  # 2 深蓝
    RGBColor(0x5C, 0x7A, 0x3A),  # 3 墨绿
    RGBColor(0x8B, 0x4F, 0xA8),  # 4 暗紫
    RGBColor(0xC4, 0x5A, 0x3C),  # 5 赭红
    RGBColor(0x8B, 0x2E, 0x2E),  # 6 深红
    RGBColor(0xC4, 0x9A, 0x3C),  # 7 金色
]

# ============================================================
# 辅助函数
# ============================================================
def set_slide_bg(slide, color):
    ""“设置幻灯片纯色背景”""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, 
                 font_size=14, color=C_TEXT_BODY, bold=False, 
                 alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei',
                 anchor=MSO_ANCHOR.TOP):
    ""“添加文本框”""
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    # 设置东亚字体
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), font_name)
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, 
                      font_size=12, color=C_TEXT_BODY, bold_first=False,
                      font_name='Microsoft YaHei', line_spacing=1.3,
                      anchor=MSO_ANCHOR.TOP):
    ""“添加多行文本框”""
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, c, b = line_data, font_size, color, (bold_first and i == 0)
        elif isinstance(line_data, tuple):
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else font_size
            c = line_data[2] if len(line_data) > 2 else color
            b = line_data[3] if len(line_data) > 3 else (bold_first and i == 0)
        else:
            continue
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
        p.space_after = Pt(2)
        
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn('a:eaTypeface'), font_name)
    
    return txBox

def add_shape_rect(slide, left, top, width, height, fill_color, 
                   border_color=None, corner_radius=None):
    ""“添加矩形”""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_line(slide, left, top, width, color=C_MID_GRAY, thickness=1):
    ""“添加横线”""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(left), Cm(top), Cm(width), Cm(thickness / 25.4 * 0.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_page_number(slide, page_num, total=14):
    ""“添加页码”""
    add_text_box(slide, 22.5, 18.2, 3, 0.6,
                 f“七擒孟获 · 游戏介绍  {page_num:02d} / {total:02d}",
                 font_size=7, color=C_MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, num, title):
    ""“添加章节标题区”""
    # 左侧色条
    add_shape_rect(slide, 0, 0, 0.3, 19.05, C_PRIMARY)
    # 章节号
    add_text_box(slide, 1.2, 0.8, 3, 1.0, f"{num:02d}",
                 font_size=42, color=C_PRIMARY, bold=True, font_name='Georgia')
    # 章节标题
    add_text_box(slide, 3.5, 0.8, 15, 1.0, title,
                 font_size=26, color=C_TEXT_DARK, bold=True)
    # 分割线
    add_line(slide, 1.2, 1.9, 22, C_PRIMARY, 3)

def add_subtitle_line(slide, text, top=2.2):
    ""“添加副标题/观点句”""
    add_text_box(slide, 1.2, top, 22.5, 0.7, text,
                 font_size=13, color=C_DARK_GRAY, bold=False)

# ============================================================
# 创建演示文稿
# ============================================================
prs = Presentation()
prs.slide_width = Cm(25.4)   # 16:10
prs.slide_height = Cm(19.05)

# ============================================================
# Page 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
set_slide_bg(slide, C_DARK)

# 顶部装饰线
add_shape_rect(slide, 0, 0, 25.4, 0.15, C_PRIMARY)

# 左侧大面积朱砂红色块
add_shape_rect(slide, 0, 0, 7.5, 19.05, C_PRIMARY)

# 主标题
add_text_box(slide, 1.2, 3.5, 5.5, 2.5, “七擒\n孟获”,
             font_size=52, color=C_WHITE, bold=True, font_name='SimHei')

# 副标题
add_text_box(slide, 1.2, 9.5, 5.5, 1.2, “像素三国\n横版动作叙事游戏”,
             font_size=16, color=RGBColor(0xFF, 0xE0, 0xE0), font_name='Microsoft YaHei')

# 左侧底部装饰
add_text_box(slide, 1.2, 15.5, 5.5, 2.0, “打败 Boss 不是终点，\n而是故事的开始。”,
             font_size=12, color=RGBColor(0xFF, 0xC0, 0xC0), font_name='Microsoft YaHei')

# 右侧信息区
add_multiline_box(slide, 9.5, 6.0, 14, 6, [
    (“腾讯云黑客松 · 叙事类游戏赛道”, 14, C_MID_GRAY),
    ("", 8, C_MID_GRAY),
    (“横版射击 × 七擒七纵 × 攻心为上”, 18, C_ACCENT_GOLD, True),
    ("", 10, C_MID_GRAY),
    ("Phase 2 MVP — Phaser 3 + TypeScript + Vite", 12, C_MID_GRAY),
    ("Web 可玩 Demo，浏览器直跑”, 12, C_MID_GRAY),
    ("", 10, C_MID_GRAY),
    (“技术栈：Phaser 3 · TypeScript · CloudBase · 腾讯混元 AI", 10, RGBColor(0x6B, 0x72, 0x80)),
], font_size=14, color=C_LIGHT_GRAY)

# 底部版本信息
add_text_box(slide, 9.5, 17.0, 14, 0.6,
             "v1.1 · 2026.06 · 像素三国视觉概念”,
             font_size=8, color=C_MID_GRAY)

# 底部装饰线
add_shape_rect(slide, 0, 18.9, 25.4, 0.15, C_PRIMARY)

# ============================================================
# Page 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 0, “目录”)

toc_items = [
    ("01", “作品定位”, “核心价值主张”),
    ("02", “故事背景”, “公元225年南中之战”),
    ("03", “核心循环”, “擒而不杀的游戏机制”),
    ("04", “七关概览”, “七擒七纵关卡设计”),
    ("05", “心态弧线”, “孟获七阶段变化”),
    ("06", “角色弧线”, “三条成长线交织”),
    ("07", “叙事创新”, “与传统动作游戏的对比”),
    ("08", “美术方向”, “像素三国视觉体系”),
    ("09", “技术方案”, "Phaser 3 + AI 增强”),
    ("10", "MVP 演示”, “第1关完整纵切”),
    ("11", “团队与展望”, “下一步计划”),
    ("12", “结尾”, “核心理念与 Q&A"),
]

for i, (num, title, desc) in enumerate(toc_items):
    row = i // 2
    col = i % 2
    x = 1.2 + col * 12.0
    y = 3.2 + row * 2.4
    
    # 编号
    add_text_box(slide, x, y, 1.5, 1.0, num,
                 font_size=28, color=C_PRIMARY, bold=True, font_name='Georgia')
    # 标题
    add_text_box(slide, x + 2.0, y, 8, 0.8, title,
                 font_size=15, color=C_TEXT_DARK, bold=True)
    # 描述
    add_text_box(slide, x + 2.0, y + 0.7, 8, 0.6, desc,
                 font_size=9, color=C_MID_GRAY)


# ============================================================
# Page 3: 作品定位
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 1, “作品定位”)
add_subtitle_line(slide, “动作游戏的胜利，被改写成叙事起点”)

# 一句话描述框
add_shape_rect(slide, 1.2, 3.0, 22.5, 2.5, C_WHITE, C_PRIMARY, True)

add_multiline_box(slide, 1.8, 3.2, 21.3, 2.2, [
    (“核心定义”, 12, C_PRIMARY, True),
    (“一款像素三国风格的横版动作游戏，以“擒而不杀、纵而再战”的七擒七纵循环为叙事引擎，”,
     14, C_TEXT_DARK, False),
    (“让玩家扮演汉军无名士卒，在诸葛亮的计谋指引下，经历七场与蛮王孟获的宿命对决。”,
     14, C_TEXT_DARK, False),
])

# 四大体验维度
dimensions = [
    (“战斗感”, “魂斗罗式横版射击\n爽快直接，像素子弹横飞”, "🎮"),
    (“叙事感”, “七擒七纵循环结构\n每场Boss战结束才是故事开始”, "📖"),
    (“文化感”, “三国智慧+蛮夷文化\n立体呈现，非标签式古装”, "🏯"),
    (“成长感”, “孟获从傲慢到折服\n玩家亲身见证心态弧线”, "🌱"),
]

for i, (title, desc, icon) in enumerate(dimensions):
    x = 1.2 + i * 5.9
    y = 6.5
    
    # 卡片背景
    add_shape_rect(slide, x, y, 5.4, 4.5, C_WHITE, None, True)
    
    # 标题
    add_text_box(slide, x + 0.4, y + 0.3, 4.6, 0.8, title,
                 font_size=18, color=C_PRIMARY, bold=True)
    # 描述
    add_text_box(slide, x + 0.4, y + 1.3, 4.6, 2.5, desc,
                 font_size=12, color=C_TEXT_BODY)

# 底部差异化定位
add_shape_rect(slide, 1.2, 12.0, 22.5, 1.5, C_DARK)
add_multiline_box(slide, 1.8, 12.2, 21.3, 1.2, [
    (“差异化定位”, 11, C_ACCENT_GOLD, True),
    ("\u201c大多数动作游戏的 Boss 打倒即结束。《七擒孟获》的独特性在于：打倒 Boss 只是故事的开始。\u201d", 
     12, C_LIGHT_GRAY, False),
])

# 核心机制标签
tags = [
    (“擒而不杀”, C_PRIMARY),
    (“纵而再战”, C_DEEP_BLUE),
    (“小兵视角”, C_INK_GREEN),
    ("AI计谋”, C_PURPLE),
]
for i, (tag, color) in enumerate(tags):
    x = 1.2 + i * 5.9
    add_shape_rect(slide, x, 14.2, 5.4, 1.0, color)
    add_text_box(slide, x, 14.4, 5.4, 0.6, tag,
                 font_size=13, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 1)

# ============================================================
# Page 4: 故事背景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 2, “故事背景”)
add_subtitle_line(slide, “公元225年：南中之战，真正的战场在人心”)

# 时间线
timeline_data = [
    ("223年”, “刘备白帝城托孤\n刘禅继位，诸葛亮辅政”),
    ("225年春”, “南中蛮族叛乱\n诸葛亮率军南征”),
    ("225年夏”, “进入南中腹地\n七擒七纵交锋”),
    ("225年秋”, “孟获折服归降\n南中平定”),
    ("225年冬”, “班师回成都\n奠定北伐基础”),
]

for i, (year, event) in enumerate(timeline_data):
    x = 1.2 + i * 4.6
    y = 3.0
    
    # 年份圆点
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x + 1.2), Cm(y), Cm(1.6), Cm(1.6))
    dot.fill.solid()
    dot.fill.fore_color.rgb = LEVEL_COLORS[i] if i < 7 else C_PRIMARY
    dot.line.fill.background()
    
    # 年份
    add_text_box(slide, x + 0.2, y + 0.1, 3.5, 0.6, year,
                 font_size=13, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, y + 0.6, 3.5, 0.6, str(i+1),
                 font_size=10, color=RGBColor(0xFF,0xFF,0xFF), alignment=PP_ALIGN.CENTER)
    
    # 连接线
    if i < 4:
        add_shape_rect(slide, x + 3.0, y + 0.75, 1.8, 0.06, LEVEL_COLORS[i])
    
    # 事件描述
    add_text_box(slide, x, y + 2.2, 3.8, 2.0, event,
                 font_size=10, color=C_TEXT_BODY, alignment=PP_ALIGN.CENTER)

# 史书记载
add_shape_rect(slide, 1.2, 6.5, 22.5, 3.8, C_DARK)
add_multiline_box(slide, 1.8, 6.7, 21.3, 3.5, [
    (“史书记载 ——《三国志·蜀书·诸葛亮传》裴松之注引《汉晋春秋》”, 11, C_ACCENT_GOLD, True),
    ("", 6, C_LIGHT_GRAY),
    ("\u201c亮至南中，所在战捷。闻孟获者，为夷、汉所服，募生致之。\u201d", 13, C_LIGHT_GRAY, False),
    (“既得，使观于营阵之间，问曰：\u2018此军何如？\u2019", 13, C_LIGHT_GRAY, False),
    (“获对曰：\u2018向者不知虚实，故败。\u2019亮笑，纵使更战。”, 13, C_LIGHT_GRAY, False),
    (“七纵七禽，而亮犹遣获。获止不去，曰：\u2018公，天威也，南人不复反矣。\u2019", 13, C_LIGHT_GRAY, False),
])

# 创作原则
add_multiline_box(slide, 1.2, 11.2, 22.5, 3.0, [
    (“历史还原度取舍原则”, 12, C_PRIMARY, True),
    (“核心事件严格遵循 · 具体战斗艺术加工 · 蛮夷文化立体呈现 · 人物基于史实延伸 · 台词整体创作”, 
     10, C_TEXT_BODY, False),
])

# 地理示意
add_shape_rect(slide, 1.2, 14.5, 22.5, 1.5, C_WHITE)
add_multiline_box(slide, 1.6, 14.6, 21.8, 1.2, [
    (“游戏舞台：南中地理”, 11, C_INK_GREEN, True),
    (“西洱河（平原）→ 泸水（天险）→ 紫荆山（密林）→ 秃龙洞（洞窟）→ 蛮族营寨 → 要塞 → 银坑洞王庭”, 
     10, C_TEXT_BODY, False),
])

add_page_number(slide, 2)

# ============================================================
# Page 5: 核心循环
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 3, “核心循环”)
add_subtitle_line(slide, “擒而不杀，纵而再战 \u2014 每一次\u201c打赢\u201d都是新故事的开始”)

# 游戏循环图
loop_items = [
    (“开场叙事”, “诸葛亮授计\nAI动态台词”, C_PRIMARY),
    (“关卡战斗”, “横版射击\n中途叙事对话”, C_DEEP_BLUE),
    ("Boss 战”, “孟获登场\n独特战斗机制”, C_INK_GREEN),
    (“擒获叙事”, “孟获反应\n释放抉择”, C_PURPLE),
    (“关卡结算”, “评价系统\n解锁碎片”, RGBColor(0xC4, 0x5A, 0x3C)),
]

for i, (step, desc, color) in enumerate(loop_items):
    x = 1.2 + i * 4.8
    y = 3.2
    
    # 步骤卡片
    add_shape_rect(slide, x, y, 4.3, 3.8, C_WHITE, None, True)
    
    # 步骤编号
    add_shape_rect(slide, x + 1.2, y - 0.5, 1.9, 1.0, color)
    add_text_box(slide, x + 1.2, y - 0.4, 1.9, 0.6, f"0{i+1}",
                 font_size=16, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 步骤名
    add_text_box(slide, x + 0.3, y + 0.8, 3.7, 0.8, step,
                 font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 描述
    add_text_box(slide, x + 0.3, y + 1.8, 3.7, 1.5, desc,
                 font_size=11, color=C_TEXT_BODY, alignment=PP_ALIGN.CENTER)
    
    # 箭头
    if i < 4:
        add_text_box(slide, x + 3.8, y + 1.3, 1.2, 1.0, "→",
                     font_size=28, color=C_MID_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# 循环说明
add_shape_rect(slide, 1.2, 8.0, 22.5, 2.2, C_DARK)
add_multiline_box(slide, 1.8, 8.2, 21.3, 1.8, [
    (“核心循环 × 7 次”, 13, C_ACCENT_GOLD, True),
    ("[主菜单] → [选关/继续] → 开场叙事 → 关卡战斗 → Boss战 → 擒获叙事 → 关卡结算 → 返回选关”,
     12, C_LIGHT_GRAY, False),
    (“循环 7 次后 → 终章·第七擒·折服 → 结局过场 → 制作名单”, 12, C_LIGHT_GRAY, False),
])

# 叙事介入方式
add_shape_rect(slide, 1.2, 11.0, 10.5, 4.5, C_WHITE)
add_multiline_box(slide, 1.6, 11.2, 9.8, 4.0, [
    (“叙事介入方式”, 13, C_PRIMARY, True),
    (“图文过场（关键节点）: 每关开场/擒获后”, 10, C_TEXT_BODY),
    (“对话框（过程叙述）: 关卡中途/Boss战前”, 10, C_TEXT_BODY),
    (“时长设计: 过场15-60秒 / 对话框2-3秒”, 10, C_TEXT_BODY),
    (“可跳过性: 除首通擒获过场外，均可跳过”, 10, C_TEXT_BODY),
    ("", 6, C_TEXT_BODY),
    (“武器渐进: 基础弩→连弩(第3关)→火箭(第5关)→诸葛连弩(第7关)", 10, C_PRIMARY, True),
])

# 战斗系统
add_shape_rect(slide, 13.0, 11.0, 10.7, 4.5, C_WHITE)
add_multiline_box(slide, 13.4, 11.2, 10.0, 4.0, [
    (“战斗操作（魂斗罗式）”, 13, C_DEEP_BLUE, True),
    (“移动：左右横移，标准平台移动”, 10, C_TEXT_BODY),
    (“跳跃：单段跳 + 空中微调”, 10, C_TEXT_BODY),
    (“射击：八方向射击（上下左右+斜向）”, 10, C_TEXT_BODY),
    (“闪避：短无敌帧翻滚（第2关解锁）”, 10, C_TEXT_BODY),
    ("", 6, C_TEXT_BODY),
    (“故事碎片：每关隐藏1块，集齐7块解锁完整史书记载”, 10, C_INK_GREEN, True),
])

add_page_number(slide, 3)

# ============================================================
# Page 6: 七关概览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 4, “七关概览”)
add_subtitle_line(slide, “每关机制变化反映孟获心态变化 — 关卡即叙事”)

levels = [
    (“第1擒”, “西洱河初战”, “平原”, “翠绿+天蓝”, "★", “基础教学：移动/跳跃/射击”, “孟获骑象，轻敌冒进”),
    (“第2擒”, “泸水天险”, “河岸”, “深蓝+灰白”, "★★", “跳跃挑战：水流+浮水平台”, “乘船水战，利用地形”),
    (“第3擒”, “紫荆山伏兵”, “密林”, “墨绿+褐黄”, "★★★", “视野遮挡+伏击+隐身识别”, “丛林猎手，设伏偷袭”),
    (“第4擒”, “秃龙洞毒雾”, “洞窟”, “暗紫+荧光绿”, "★★★★", “环境伤害+视野黑暗+资源管理”, “联合祝融夫人，毒雾”),
    (“第5擒”, “蛮族联军”, “营寨”, “赭红+土黄”, "★★★★", “多波次敌兵+双Boss战”, “倾巢而出，赌上一切”),
    (“第6擒”, “木兽炎阵”, “要塞”, “深红+铁灰”, "★★★★★", “火焰陷阱+木兽机关+反制”, “出奇制胜，用尽底牌”),
    (“第7擒”, “银坑洞决战”, “王庭”, “金+深棕”, "★★★★★", “全机制综合+三阶段Boss", “背水一战，最终折服”),
]

# 表头
headers = [“关卡”, “名称”, “场景”, “色调”, “难度”, “核心机制”, "Boss战特色”]
col_widths = [2.2, 3.0, 1.8, 2.8, 1.5, 5.0, 6.0]
x_start = 1.2
y_start = 3.0

# 表头背景
cum_x = x_start
for j, (header, cw) in enumerate(zip(headers, col_widths)):
    add_shape_rect(slide, cum_x, y_start, cw, 1.0, C_DARK)
    add_text_box(slide, cum_x + 0.15, y_start + 0.15, cw - 0.3, 0.7, header,
                 font_size=9, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    cum_x += cw

# 表格行
for i, (level, name, scene, tone, diff, mech, boss) in enumerate(levels):
    y = y_start + 1.0 + i * 1.4
    bg = C_WHITE if i % 2 == 0 else RGBColor(0xF8, 0xF7, 0xF3)
    
    cum_x = x_start
    texts = [level, name, scene, tone, diff, mech, boss]
    sizes = [10, 11, 9, 8, 9, 8, 8]
    for j, (text, cw, size) in enumerate(zip(texts, col_widths, sizes)):
        # 行背景
        add_shape_rect(slide, cum_x, y, cw, 1.3, bg)
        
        # 左侧装饰色条（关卡色调）
        if j == 0:
            add_shape_rect(slide, cum_x, y, 0.2, 1.3, LEVEL_COLORS[i])
        
        color = C_TEXT_DARK if j <= 1 else C_TEXT_BODY
        bold = j <= 1
        add_text_box(slide, cum_x + 0.3, y + 0.1, cw - 0.5, 1.1, text,
                     font_size=size, color=color, bold=bold, 
                     alignment=PP_ALIGN.CENTER if j < 4 else PP_ALIGN.LEFT)
        cum_x += cw

add_page_number(slide, 4)

# ============================================================
# Page 7: 孟获心态弧线
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 5, “孟获心态弧线”)
add_subtitle_line(slide, “从傲慢到折服 — 七擒七纵的本质是让对手\u201c被理解\u201d，而非\u201c被击败\u201d")

moods = [
    (“第1擒\n傲慢”, "\u201c诡计！汉人狡诈！\n我不服！\u201d", LEVEL_COLORS[0]),
    (“第2擒\n辩解”, "\u201c此乃地利之失，\n非战之罪！\u201d", LEVEL_COLORS[1]),
    (“第3擒\n困惑”, "\u201c……汝究竟\n欲如何？\u201d", LEVEL_COLORS[2]),
    (“第4擒\n倔强”, "\u201c纵我七次，\n我也不降！\u201d", LEVEL_COLORS[3]),
    (“第5擒\n挣扎”, "\u201c我蛮族男儿，\n岂能降汉！\u201d", LEVEL_COLORS[4]),
    (“第6擒\n反思”, "\u201c你是妖魔吗？\n你们来做什么？\u201d", LEVEL_COLORS[5]),
    (“第7擒\n折服”, "\u201c公，天威也。\n南人不复反矣。\u201d", LEVEL_COLORS[6]),
]

for i, (stage, quote, color) in enumerate(moods):
    x = 1.2 + i * 3.4
    y = 3.2
    
    # 阶段卡片
    add_shape_rect(slide, x, y, 3.1, 7.5, C_WHITE, None, True)
    
    # 顶部色条
    add_shape_rect(slide, x, y, 3.1, 0.5, color)
    
    # 阶段名
    add_text_box(slide, x, y + 0.7, 3.1, 1.5, stage,
                 font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 箭头向下
    add_text_box(slide, x + 1.0, y + 2.5, 1.0, 0.6, "↓",
                 font_size=16, color=color, alignment=PP_ALIGN.CENTER)
    
    # 台词
    add_text_box(slide, x + 0.2, y + 3.2, 2.7, 3.0, quote,
                 font_size=10, color=C_TEXT_BODY, alignment=PP_ALIGN.CENTER)

# 底部核心洞察
add_shape_rect(slide, 1.2, 11.5, 22.5, 3.5, C_DARK)
add_multiline_box(slide, 1.8, 11.7, 21.3, 3.0, [
    (“设计洞察”, 12, C_ACCENT_GOLD, True),
    ("", 6, C_LIGHT_GRAY),
    (“孟获反抗的根本原因不是好战，而是恐惧 —— 恐惧南中失去独立，恐惧族人被奴役。”, 13, C_LIGHT_GRAY, False),
    (“诸葛亮最终降服的，不是孟获的武力，而是他的恐惧。”, 13, C_LIGHT_GRAY, True),
    ("", 6, C_LIGHT_GRAY),
    (“核心叙事钩子：为什么打赢了还要放人？→ 答案随7关逐步揭晓。”, 11, RGBColor(0x6B, 0x72, 0x80), False),
])

add_page_number(slide, 5)

# ============================================================
# Page 8: 三条角色弧线
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 6, “三条角色弧线”)
add_subtitle_line(slide, “诸葛亮、孟获、玩家小卒共同完成“心战”的闭环”)

# 三条角色线
characters = [
    {
        "title": “玩家小卒”,
        "subtitle": “汉军无名士卒 — 玩家化身”,
        "color": C_PRIMARY,
        "desc": “一个普通蜀中青年，应征入伍。\n他不是名将之后，没有显赫家世——\n就是在成都街巷里长大的普通人。\n被诸葛亮选为前锋斥候，亲历\n七擒孟获的每一个关键时刻。”,
        "arc": “困惑 → 思考 → 共情 → 见证”,
        "quote": "“丞相为何放走敌人？\n→ 好像有点理解了\n→ 孟获也是个有血有肉的人\n→ 我见证了历史”",
    },
    {
        "title": “诸葛亮”,
        "subtitle": “蜀汉丞相 · 南征统帅”,
        "color": C_DEEP_BLUE,
        "desc": "45岁。羽扇纶巾，深蓝色长袍。\n不是一个全知全能的\u201c神仙\u201d，\n而是一个有智慧、有远见、\n有温度的人。善于观察和分析，\n但也会面临不确定性。”,
        "arc": “战略导师 \u00b7 不神化\n以耐心换和平”,
        "quote": "\u201c攻心为上，攻城为下。\n心战为上，兵战为下。\u201d\n\n相信人心可以被感化。\n真正的胜利不是杀戮，\n是让对方心甘情愿。”,
    },
    {
        "title": “孟获”,
        "subtitle": “南中蛮族联盟领袖 · Boss",
        "color": C_INK_GREEN,
        "desc": “约40岁。虎背熊腰，赤膊纹身，\n头戴兽骨盔，手持战斧。\n不是一个被\u201c打败\u201d的Boss，\n而是一个被\u201c理解\u201d的人。\n勇武且有自尊，有部落荣誉感。”,
        "arc": “对手 \u2192 宿敌 \u2192 伙伴\n从抗拒到理解到归心”,
        "quote": "\u201c我不是不相信他。\n我只是不相信汉人。\u201d\n\n最终说出：\n\u201c公，天威也。\n南人不复反矣。\u201d",
    },
]

for i, char in enumerate(characters):
    x = 1.2 + i * 8.0
    y = 3.0
    
    # 角色卡片
    add_shape_rect(slide, x, y, 7.5, 7.5, C_WHITE, None, True)
    
    # 顶部色条
    add_shape_rect(slide, x, y, 7.5, 0.4, char["color"])
    
    # 角色名
    add_text_box(slide, x + 0.5, y + 0.6, 6.5, 0.8, char["title"],
                 font_size=20, color=char["color"], bold=True)
    add_text_box(slide, x + 0.5, y + 1.3, 6.5, 0.5, char["subtitle"],
                 font_size=9, color=C_MID_GRAY)
    
    # 描述
    add_text_box(slide, x + 0.5, y + 2.0, 6.5, 3.0, char["desc"],
                 font_size=9, color=C_TEXT_BODY)
    
    # 弧线
    add_shape_rect(slide, x + 0.5, y + 4.8, 6.5, 0.8, RGBColor(0xF0, 0xEF, 0xEA))
    add_text_box(slide, x + 0.7, y + 4.9, 6.1, 0.6, char["arc"],
                 font_size=9, color=char["color"], bold=True)
    
    # 金句
    add_text_box(slide, x + 0.5, y + 5.8, 6.5, 1.5, char["quote"],
                 font_size=8, color=C_DARK_GRAY)

# 底部关系图
add_shape_rect(slide, 1.2, 11.5, 22.5, 1.8, C_DARK)
add_multiline_box(slide, 1.8, 11.7, 21.3, 1.4, [
    (“三条线交织 — 共同完成“心战”的闭环”, 12, C_ACCENT_GOLD, True),
    (“诸葛亮（授计者）→ 玩家小卒（执行者/见证者）→ 孟获（被感化者）”, 10, C_LIGHT_GRAY, False),
    (“师徒式关系：诸葛亮对玩家的解释，也是在对所有玩家“授课”。第7关那句“你可明白了？”同时问向角色与玩家。”, 10, C_LIGHT_GRAY, False),
])

add_page_number(slide, 6)

# ============================================================
# Page 9: 叙事创新
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 7, “叙事创新”)
add_subtitle_line(slide, “把”攻心为上”从台词，变成可见、可玩、可验证的机制”)

# 对比
# 传统动作游戏
add_shape_rect(slide, 1.2, 3.2, 10.5, 5.0, RGBColor(0xF0, 0xEF, 0xEA))
add_text_box(slide, 1.8, 3.4, 9.3, 0.6, “传统动作游戏”,
             font_size=16, color=C_MID_GRAY, bold=True)
add_multiline_box(slide, 1.8, 4.2, 9.3, 3.5, [
    ("Boss 被打倒”, 14, C_MID_GRAY, True),
    ("↓", 14, C_MID_GRAY),
    (“死亡 / 爆炸 / 通关画面”, 14, C_MID_GRAY, True),
    ("↓", 14, C_MID_GRAY),
    (“下一关”, 14, C_MID_GRAY, True),
    ("", 8, C_MID_GRAY),
    (“叙事 = 战斗的句号”, 12, C_DARK_GRAY, True),
])

# VS
add_text_box(slide, 10.5, 5.0, 4.4, 1.2, "VS",
             font_size=40, color=C_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# 七擒孟获
add_shape_rect(slide, 13.0, 3.2, 10.7, 5.0, C_DARK)
add_text_box(slide, 13.6, 3.4, 9.5, 0.6, “《七擒孟获》”,
             font_size=16, color=C_ACCENT_GOLD, bold=True)
add_multiline_box(slide, 13.6, 4.2, 9.5, 3.5, [
    ("Boss 被打倒”, 14, C_ACCENT_GOLD, True),
    ("↓", 14, C_ACCENT_GOLD),
    (“被擒获 · 诸葛亮释放”, 14, C_ACCENT_GOLD, True),
    ("↓", 14, C_ACCENT_GOLD),
    (“孟获带新理由再战”, 14, C_ACCENT_GOLD, True),
    ("", 8, C_ACCENT_GOLD),
    (“叙事 = 新故事的开始”, 12, C_PRIMARY, True),
])

# 叙事创新点
innovations = [
    (“擒而不杀”, “每次战胜孟获后触发释放叙事节点。诸葛亮下令松绑，孟获震惊、困惑、反思——这些情绪变化成为下一步关卡设计的动机。”),
    (“纵而再战”, “下次相遇时孟获更强、更狡猾，关卡机制同步深化。第1关轻敌→第7关全力以赴，难度和复杂度与心态弧线同步增长。”),
    (“玩家追问”, “首通不可跳过的擒获过场确保玩家亲历“为什么放他走”的困惑。验收标准：玩家通关后能主动说出——“我打赢了，但故事还没结束。”。”),
    ("AI计谋增强”, “腾讯混元大模型根据玩家上关行为（用时/死亡/受伤）动态生成诸葛亮计谋台词，有本地 fallback 保证不阻塞主流程。”),
]

for i, (title, desc) in enumerate(innovations):
    x = 1.2 + (i % 2) * 12.0
    y = 8.8 + (i // 2) * 2.8
    
    # 卡片
    add_shape_rect(slide, x, y, 11.0, 2.5, C_WHITE, None, True)
    
    # 标题
    add_text_box(slide, x + 0.4, y + 0.2, 10.2, 0.6, title,
                 font_size=14, color=C_PRIMARY if i < 2 else C_DEEP_BLUE, bold=True)
    # 描述
    add_text_box(slide, x + 0.4, y + 0.9, 10.2, 1.4, desc,
                 font_size=9, color=C_TEXT_BODY)

add_page_number(slide, 7)

# ============================================================
# Page 10: 美术方向
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 8, “美术方向：像素三国”)
add_subtitle_line(slide, “复古街机的动作感 + 国风水墨的历史质感”)

# 核心风格
style_items = [
    (“像素基底”, "16×16 或 32×32 Sprite，保持魂斗罗式的爽快感。角色与战斗保持清晰可读。”, C_PRIMARY),
    (“国风点缀”, "UI 边框：水墨笔触纹理\n标题/关卡名：毛笔手写字体\n过场插画：水墨渲染 + 像素角色\n粒子效果：墨滴飞溅替代传统火花”, C_INK_GREEN),
    (“低饱和调色”, “朱砂红、墨绿、古铜、深蓝为主色调，避免高饱和“游戏感”。营造历史的厚重与沉淀。”, C_DEEP_BLUE),
    (“色调叙事”, “七关从翠绿到金棕，环境色调映射孟获心态从轻浮到沉淀的变化。”, C_PURPLE),
]

for i, (title, desc, color) in enumerate(style_items):
    y = 3.0 + i * 3.0
    x = 1.2
    
    # 色条
    add_shape_rect(slide, x, y, 0.3, 2.5, color)
    
    # 标题
    add_text_box(slide, x + 0.8, y + 0.1, 3.5, 0.6, title,
                 font_size=16, color=color, bold=True)
    # 描述
    add_text_box(slide, x + 0.8, y + 0.8, 10.0, 1.6, desc,
                 font_size=10, color=C_TEXT_BODY)

# 配色展示
add_shape_rect(slide, 13.0, 3.0, 11.0, 12.0, C_DARK)
add_text_box(slide, 13.5, 3.3, 10.0, 0.6, “七关环境色调”,
             font_size=13, color=C_ACCENT_GOLD, bold=True)

level_tones = [
    (“第1擒”, “西洱河平原”, “翠绿+天蓝”, LEVEL_COLORS[0], “开阔、初战、轻敌”),
    (“第2擒”, “泸水河岸”, “深蓝+灰白”, LEVEL_COLORS[1], “险峻、湍急、沉着”),
    (“第3擒”, “紫荆山密林”, “墨绿+褐黄”, LEVEL_COLORS[2], “压抑、伏击、智斗”),
    (“第4擒”, “秃龙洞洞窟”, “暗紫+荧光绿”, LEVEL_COLORS[3], “诡异、毒雾、坚持”),
    (“第5擒”, “蛮族联盟营寨”, “赭红+土黄”, LEVEL_COLORS[4], “宏大、联战、暗流”),
    (“第6擒”, “蛮夷要塞”, “深红+铁灰”, LEVEL_COLORS[5], “烈火、机关、反制”),
    (“第7擒”, “银坑洞王庭”, “金+深棕”, LEVEL_COLORS[6], “壮阔、决战、折服”),
]

for i, (level, scene, tone, color, mood) in enumerate(level_tones):
    y = 4.2 + i * 1.5
    
    # 色块
    add_shape_rect(slide, 13.5, y, 0.8, 1.2, color)
    
    # 关卡名
    add_text_box(slide, 14.5, y + 0.05, 2.5, 0.5, level,
                 font_size=10, color=C_LIGHT_GRAY, bold=True)
    add_text_box(slide, 14.5, y + 0.55, 2.5, 0.4, scene,
                 font_size=8, color=C_MID_GRAY)
    
    # 色调描述
    add_text_box(slide, 17.0, y + 0.05, 2.5, 0.5, tone,
                 font_size=9, color=C_ACCENT_GOLD)
    add_text_box(slide, 17.0, y + 0.55, 6.5, 0.4, mood,
                 font_size=8, color=C_MID_GRAY)

add_page_number(slide, 8)

# ============================================================
# Page 11: 技术方案
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 9, “技术方案”)
add_subtitle_line(slide, "Web 可玩 Demo：Phaser 3 + TypeScript + Vite + CloudBase")

# 技术栈卡片
tech_stack = [
    ("Phaser 3.60+", “游戏引擎”, “零安装，浏览器直跑\n评审官零门槛\n2D 游戏框架成熟稳定”, C_PRIMARY),
    ("TypeScript 5.0+", “编程语言”, “类型安全，大型项目可维护\n代码提示友好\n重构安全”, C_DEEP_BLUE),
    ("Vite 5.0+", “构建工具”, “极速 HMR，开发体验好\n秒级冷启动\n天然支持 TypeScript", C_INK_GREEN),
    ("CloudBase", “部署平台”, “腾讯云静态托管\n一键部署\n腾讯生态深度集成”, C_PURPLE),
]

for i, (name, role, desc, color) in enumerate(tech_stack):
    x = 1.2 + i * 6.0
    y = 3.2
    
    add_shape_rect(slide, x, y, 5.5, 4.5, C_WHITE, None, True)
    add_shape_rect(slide, x, y, 5.5, 0.4, color)
    
    add_text_box(slide, x + 0.3, y + 0.6, 4.9, 0.8, name,
                 font_size=18, color=color, bold=True)
    add_text_box(slide, x + 0.3, y + 1.3, 4.9, 0.5, role,
                 font_size=10, color=C_MID_GRAY)
    add_text_box(slide, x + 0.3, y + 2.0, 4.9, 2.0, desc,
                 font_size=10, color=C_TEXT_BODY)

# AI 系统
add_shape_rect(slide, 1.2, 8.5, 22.5, 4.5, C_DARK)
add_multiline_box(slide, 1.8, 8.7, 21.3, 4.0, [
    ("AI 动态计谋系统（加分项）”, 13, C_ACCENT_GOLD, True),
    ("", 6, C_LIGHT_GRAY),
    (“流程：玩家上关数据（用时/死亡/受伤）→ 行为分析（激进/稳健/快/慢）→ 腾讯混元大模型 → 诸葛亮计谋台词”, 
     11, C_LIGHT_GRAY, False),
    ("", 6, C_LIGHT_GRAY),
    (“设计原则”, 11, C_PRIMARY, True),
    ("· AI 只生成非关键短建议（战术/心性提醒），主线剧情由人类写定”, 10, C_LIGHT_GRAY),
    ("· 本地 fallback 模板优先，不允许 API 阻塞主流程”, 10, C_LIGHT_GRAY),
    ("· 诸葛亮语气约束：战前严肃简洁 / 日常温和有洞察 / 关键节点厚重有分量”, 10, C_LIGHT_GRAY),
    ("", 6, C_LIGHT_GRAY),
    (“架构：5层分层（Scenes·UI·Systems·Entities·Core/Utils），场景即入口，System无状态，Entity自包含，Data驱动”, 
     10, RGBColor(0x6B, 0x72, 0x80), False),
])

add_page_number(slide, 9)

# ============================================================
# Page 12: MVP 演示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 10, "MVP 演示范围”)
add_subtitle_line(slide, “先做第1关”西洱河初战”完整纵切，证明核心体验”)

# 8分钟脚本
script_steps = [
    ("01", “诸葛亮授计”, "“此战所求，非杀戮，乃归心。我要活的。”", C_PRIMARY),
    ("02", “基础教学”, “玩家学习移动、跳跃、射击（魂斗罗式八方向）”, C_DEEP_BLUE),
    ("03", “关卡推进”, “平原行军→河岸遭遇战，蛮兵增多”, C_INK_GREEN),
    ("04", "Boss 战”, “孟获骑战象登场：象鼻横扫/踏地震击/冲锋”, C_PURPLE),
    ("05", “击败 Boss", “孟获台词：'什么？！这不可能——' 玩家以为自己通关”, RGBColor(0xC4, 0x5A, 0x3C)),
    ("06", “擒获释放”, “诸葛亮下令松绑：'将军勇则勇矣，然未识天时。'", C_PRIMARY),
    ("07", “关卡结算”, “第一擒完成，第二擒解锁。历史碎片收集。”, C_DEEP_BLUE),
]

for i, (num, title, desc, color) in enumerate(script_steps):
    x = 1.2 + (i % 4) * 5.9
    y = 3.2 + (i // 4) * 3.8
    
    # 步骤卡片
    add_shape_rect(slide, x, y, 5.4, 3.3, C_WHITE, None, True)
    
    # 编号
    add_shape_rect(slide, x + 0.3, y + 0.3, 1.2, 0.8, color)
    add_text_box(slide, x + 0.3, y + 0.35, 1.2, 0.6, num,
                 font_size=16, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 标题
    add_text_box(slide, x + 1.8, y + 0.35, 3.3, 0.6, title,
                 font_size=13, color=color, bold=True)
    # 描述
    add_text_box(slide, x + 0.3, y + 1.3, 4.8, 1.8, desc,
                 font_size=10, color=C_TEXT_BODY)

# 验收标准
add_shape_rect(slide, 1.2, 11.5, 22.5, 2.0, C_DARK)
add_multiline_box(slide, 1.8, 11.7, 21.3, 1.6, [
    (“验收标准”, 13, C_ACCENT_GOLD, True),
    (“玩家通关后能主动说出 —— “我打赢了，但故事还没结束。”", 14, C_LIGHT_GRAY, True),
    (“技术验收：30秒内学会移动+跳跃+射击 / 2分钟内遇首个叙事节点 / 5分钟内完成首次Boss战”, 10, RGBColor(0x6B, 0x72, 0x80), False),
])

add_page_number(slide, 10)

# ============================================================
# Page 13: 团队与研发计划
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_LIGHT_GRAY)
add_section_header(slide, 11, “研发进度与下一步”)
add_subtitle_line(slide, "Phase 2 MVP 已搭建完成，核心战斗可玩”)

# Phase 进度
phases = [
    ("Phase 1", "✅ 已完成”, “产品定义 & 叙事设计”, 
     "GDD / 世界观圣经 / 关卡叙事大纲 / 角色设定表 / 视觉概念定稿”, C_PRIMARY),
    ("Phase 2", "🔄 进行中”, "MVP 第1关纵切原型”,
     "5层模块架构 / Player+Enemy+Boss 实体 / 全部核心系统（含 CaptureSystem）”, C_DEEP_BLUE),
    ("Phase 3", "📋 计划中”, “第2-7关铺量”,
     “基于第1关模板批量生产 · 每关独特机制 + 叙事变体”, C_INK_GREEN),
    ("Phase 4", "📋 计划中”, "AI 集成 + 打磨”,
     “混元LLM 动态台词 / 像素美术资产 / 音效BGM / 多语言”, C_PURPLE),
]

for i, (phase, status, title, detail, color) in enumerate(phases):
    y = 3.0 + i * 3.5
    x = 1.2
    
    # 色条
    add_shape_rect(slide, x, y, 0.3, 3.0, color)
    
    # Phase + 状态
    add_text_box(slide, x + 0.8, y + 0.1, 6.0, 0.7, f"{phase}  {status}",
                 font_size=15, color=color, bold=True)
    add_text_box(slide, x + 0.8, y + 0.8, 6.0, 0.5, title,
                 font_size=12, color=C_TEXT_DARK, bold=True)
    add_text_box(slide, x + 0.8, y + 1.4, 20.0, 1.2, detail,
                 font_size=10, color=C_TEXT_BODY)

# 核心里程碑
add_shape_rect(slide, 1.2, 15.5, 22.5, 1.5, C_WHITE)
add_multiline_box(slide, 1.6, 15.6, 21.8, 1.2, [
    ("11个开发批次（Batch 0-11）· 四大通信模式解耦 · 配置化驱动 · 本地 fallback 优先”, 
     10, C_TEXT_BODY, False),
])

add_page_number(slide, 11)

# ============================================================
# Page 14: 结尾
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_DARK)

# 顶部装饰线
add_shape_rect(slide, 0, 0, 25.4, 0.15, C_PRIMARY)

# 核心金句
add_text_box(slide, 2.0, 3.5, 21.4, 3.0, “真正的胜利，\n不是消灭敌人。”,
             font_size=44, color=C_WHITE, bold=True, font_name='SimHei',
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, 2.0, 7.5, 21.4, 3.0, “而是让冲突双方\n找到共同秩序。”,
             font_size=44, color=C_ACCENT_GOLD, bold=True, font_name='SimHei',
             alignment=PP_ALIGN.CENTER)

# 分割线
add_shape_rect(slide, 8.0, 11.0, 9.4, 0.1, C_ACCENT_GOLD)

# 描述
add_text_box(slide, 3.0, 11.5, 19.4, 1.5,
             “《七擒孟获》将”七擒七纵、攻心为上”做成关卡循环、Boss 机制和玩家情绪体验。”,
             font_size=14, color=C_MID_GRAY, alignment=PP_ALIGN.CENTER)

# 标签
tags = [“横版射击”, “七擒七纵”, “小兵视角”, “攻心为上”]
for i, tag in enumerate(tags):
    x = 5.0 + i * 4.2
    add_shape_rect(slide, x, 13.5, 3.5, 1.0, C_DARK, C_ACCENT_GOLD, True)
    add_text_box(slide, x, 13.65, 3.5, 0.7, tag,
                 font_size=12, color=C_ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, 2.0, 16.0, 21.4, 1.0, "Q & A",
             font_size=28, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 2.0, 17.2, 21.4, 0.6,
             “腾讯云黑客松 · 叙事类游戏赛道 · Phaser 3 + TypeScript + Vite + CloudBase",
             font_size=9, color=C_MID_GRAY, alignment=PP_ALIGN.CENTER)

# 底部装饰线
add_shape_rect(slide, 0, 18.9, 25.4, 0.15, C_PRIMARY)

add_page_number(slide, 12)


# ============================================================
# 保存
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), '七擒孟获_游戏介绍PPT_优化版.pptx')
prs.save(output_path)
print(f"✅ PPT 生成完成！”)
print(f"📁 文件: {output_path}")
print(f"📊 共 {len(prs.slides)} 页”)
